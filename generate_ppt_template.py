#!/usr/bin/env python3
"""Generate presentation using user's template (presentation6.ppt).
Template: 10" × 7.5", Office Theme, blue title slide, white content slides.
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import copy, os
import lxml.etree as ET

# ── Paths ──
BASE = "/Users/tsy_02/llm_benchmark_reliability"
TEMPLATE = os.path.join(BASE, "presentation6.ppt")
FIG1 = os.path.join(BASE, "exp1/figures_exp1")
FIG2 = os.path.join(BASE, "exp2/figures_exp2")
FIG3 = os.path.join(BASE, "exp3/figures_exp3_shared150")

prs = Presentation(TEMPLATE)
SW = 10.0  # slide width inches
SH = 7.5   # slide height inches

# ── Colors matching template theme ──
BLUE_BG    = RGBColor(0x00, 0x42, 0x82)  # title slide background
ACCENT1    = RGBColor(0x5B, 0x9B, 0xD5)  # theme accent1 blue
ACCENT2    = RGBColor(0xED, 0x7D, 0x31)  # theme accent2 orange
ACCENT3    = RGBColor(0x70, 0xAD, 0x47)  # theme accent6 green
WHITE      = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY = RGBColor(0xE0, 0xE0, 0xE0)
MID_GRAY   = RGBColor(0x75, 0x75, 0x75)
DARK_TEXT   = RGBColor(0x44, 0x54, 0x6A)  # dk2
RED_ALERT  = RGBColor(0xE0, 0x40, 0x40)
CARD_LIGHT = RGBColor(0xF2, 0xF2, 0xF2)
CARD_BLUE  = RGBColor(0xEA, 0xF0, 0xFA)

# Layout indices from template
LY_TITLE   = 0   # Title Slide (blue bg)
LY_CONTENT = 1   # Title and Content
LY_SECTION = 2   # Section Header
LY_TWO     = 3   # Two Content
LY_COMPARE = 4   # Comparison
LY_TITLE_ONLY = 5  # Title Only + content placeholder
LY_BLANK   = 6   # Blank
LY_CAPTION = 7   # Content with Caption

# ── Remove template placeholder slides ──
# (delete slides 1-5 from template)
while len(prs.slides) > 0:
    rId = prs.slides._sldIdLst[0].get(
        '{http://schemas.openxmlformats.org/officeDocument/2006/relationships}id')
    prs.part.drop_rel(rId)
    prs.slides._sldIdLst.remove(prs.slides._sldIdLst[0])


# ── Helper functions ──
def add_slide(layout_idx):
    return prs.slides.add_slide(prs.slide_layouts[layout_idx])

def set_bg(slide, color):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_textbox(slide, left, top, width, height, text, font_size=18,
                bold=False, color=DARK_TEXT, alignment=PP_ALIGN.LEFT,
                font_name="Calibri"):
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top),
                                      Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.font.name = font_name
    p.alignment = alignment
    return tf

def add_multiline(slide, left, top, width, height, lines, font_size=14,
                  color=DARK_TEXT, bold_first=False, font_name="Calibri",
                  alignment=PP_ALIGN.LEFT):
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top),
                                      Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = line
        p.font.size = Pt(font_size)
        p.font.color.rgb = color
        p.font.name = font_name
        p.alignment = alignment
        p.space_after = Pt(font_size * 0.2)
        p.space_before = Pt(font_size * 0.05)
        if bold_first and i == 0:
            p.font.bold = True
    return tf

def add_accent_bar(slide, left, top, width=0.05, height=0.6, color=ACCENT2):
    """Add the template-style thin vertical accent bar."""
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(left), Inches(top),
        Inches(width), Inches(height))
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()

def add_hline(slide, left, top, width, color=ACCENT2, thickness=0.03):
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(left), Inches(top),
        Inches(width), Inches(thickness))
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()

def add_card(slide, left, top, width, height, bg_color=CARD_LIGHT):
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, Inches(left), Inches(top),
        Inches(width), Inches(height))
    shape.fill.solid()
    shape.fill.fore_color.rgb = bg_color
    shape.line.fill.background()
    shape.shadow.inherit = False

def add_img(slide, path, left, top, width=None, height=None):
    if os.path.exists(path):
        kw = {}
        if width: kw['width'] = Inches(width)
        if height: kw['height'] = Inches(height)
        slide.shapes.add_picture(path, Inches(left), Inches(top), **kw)

def add_table(slide, left, top, width, height, rows, cols, data,
              header_color=BLUE_BG, cell_font=11, header_font=12):
    ts = slide.shapes.add_table(rows, cols, Inches(left), Inches(top),
                                 Inches(width), Inches(height))
    table = ts.table
    for ri in range(rows):
        for ci in range(cols):
            cell = table.cell(ri, ci)
            cell.text = str(data[ri][ci])
            for p in cell.text_frame.paragraphs:
                p.font.size = Pt(header_font if ri == 0 else cell_font)
                p.font.name = "Calibri"
                p.alignment = PP_ALIGN.CENTER
                if ri == 0:
                    p.font.bold = True
                    p.font.color.rgb = WHITE
                else:
                    p.font.color.rgb = DARK_TEXT
            cf = cell.fill
            cf.solid()
            if ri == 0:
                cf.fore_color.rgb = header_color
            elif ri % 2 == 1:
                cf.fore_color.rgb = RGBColor(0xF7, 0xF7, 0xF7)
            else:
                cf.fore_color.rgb = RGBColor(0xEA, 0xF0, 0xFA)
    return table

def slide_title(slide, text, color=BLUE_BG):
    """Add standard title + accent bar to a content slide."""
    add_textbox(slide, 0.69, 0.30, 8.6, 0.65, text,
                font_size=28, bold=True, color=color, font_name="Calibri Light")
    add_accent_bar(slide, 0.64, 0.75, 0.05, 0.55, ACCENT2)


# =========================================================================
# SLIDE 1: Title (blue background)
# =========================================================================
slide = add_slide(LY_BLANK)
set_bg(slide, BLUE_BG)

add_textbox(slide, 0.6, 1.8, 8.8, 1.2,
            "How Reliable Are\nLLM Benchmark Scores?",
            font_size=38, bold=True, color=WHITE, font_name="Calibri Light",
            alignment=PP_ALIGN.LEFT)
add_accent_bar(slide, 0.45, 2.0, 0.05, 1.1, ACCENT2)

add_textbox(slide, 0.6, 3.4, 8.8, 0.7,
            "Quantifying Evaluation Noise from Prompt Wording,\nTest-Set Paraphrasing, and Item-Level Instability",
            font_size=16, color=LIGHT_GRAY, font_name="Calibri")

add_textbox(slide, 0.6, 4.5, 8.8, 0.4,
            "ST5230 Applied NLP — Group 1",
            font_size=14, color=ACCENT2, font_name="Calibri")

add_textbox(slide, 0.6, 5.4, 8.8, 0.6,
            "Benchmarks: ARC-Challenge (300) · MMLU-Pro (300)\nModels: LLaMA-3.1-8B · Qwen2.5-7B · Qwen3-32B · Qwen2.5-72B",
            font_size=12, color=RGBColor(0xA0, 0xC0, 0xE0))


# =========================================================================
# SLIDE 2: Background — The Reliability Crisis
# =========================================================================
slide = add_slide(LY_BLANK)
slide_title(slide, "Background: The Reliability Crisis")

# Left: The Problem
add_card(slide, 0.5, 1.5, 4.3, 5.5, CARD_LIGHT)
add_textbox(slide, 0.7, 1.55, 3.9, 0.35, "Benchmark Scores Are Not As They Seem",
            font_size=14, bold=True, color=RED_ALERT)

prob_lines = [
    "➊  Prompt Sensitivity",
    "    Sclar et al.: up to 76% accuracy gap",
    "    from prompt wording alone (4800 formats)",
    "",
    "➋  Ranking Instability",
    "    Alzahrani et al.: removing 0.02% of evals",
    "    can flip the #1 model",
    "",
    "➌  Test-Set Quality",
    "    Rein et al.: 3-10% label error rate in",
    "    major benchmarks; corrections shift ranking",
    "",
    "➍  Run-to-Run Variance",
    "    Pezeshkpour et al.: 10% variance even",
    "    at temperature = 0",
]
add_multiline(slide, 0.7, 2.0, 3.9, 4.5, prob_lines, font_size=11, color=DARK_TEXT)

# Right: Why It Matters
add_card(slide, 5.1, 1.5, 4.5, 5.5, CARD_BLUE)
add_textbox(slide, 5.3, 1.55, 4.1, 0.35, "Why It Matters",
            font_size=14, bold=True, color=BLUE_BG)

matter_lines = [
    "Benchmark scores drive:",
    "  • Multi-billion $ deployment decisions",
    "  • Academic claims & model comparisons",
    "  • Public leaderboard rankings (Open LLM)",
    "",
    "Yet scores are reported as single numbers",
    "without confidence intervals.",
    "",
    "  Benchmark    Known Issue         Source",
    "  ─────────  ────────────────  ────────────",
    "  MMLU          6.3% label errors   Rein et al.",
    "  ARC            Format sensitivity   Sclar et al.",
    "  HellaSwag   Ordering effects     Pezeshkpour",
    "  GSM8K        Prompt sensitivity   Mizrahi et al.",
    "",
    "→ Single-score reporting is misleading",
]
add_multiline(slide, 5.3, 2.0, 4.1, 4.5, matter_lines, font_size=11, color=DARK_TEXT)


# =========================================================================
# SLIDE 3: Related Work — Four Noise Dimensions
# =========================================================================
slide = add_slide(LY_BLANK)
slide_title(slide, "Related Work: Four Noise Dimensions")

card_data = [
    ("Prompt\nSensitivity", ACCENT2,
     ["Sclar et al. (2024):",
      "  4800+ formats tested",
      "  Up to 76% acc range",
      "Mizrahi et al. (2024):",
      "  Format > instruction",
      "Weber et al. (2023):",
      "  LLMs are poor MCQ",
      "  test-takers"]),
    ("Ranking\nStability", ACCENT1,
     ["Polo et al. (2024):",
      "  >50% reversal prob",
      "  Bootstrap CIs needed",
      "Alzahrani et al. (2024):",
      "  0.02% removal → flip",
      "Zong et al. (2024):",
      "  Score ≈ coin flip",
      "  under changed conds"]),
    ("Benchmark\nQuality", ACCENT3,
     ["Rein et al. (2023):",
      "  3-10% error rate in",
      "  4 major benchmarks",
      "Gupta et al. (2024):",
      "  Paraphrase → acc Δ",
      "Gema et al. (2024):",
      "  MMLU error analysis",
      "  and correction"]),
    ("Evaluation\nNoise",  RGBColor(0x44, 0x72, 0xC4),
     ["Pezeshkpour (2023):",
      "  10% var at temp=0",
      "  Multiple runs needed",
      "Biderman et al. (2024):",
      "  Lessons from Pythia",
      "Burnell et al. (2023):",
      "  Need CIs, not single",
      "  scores"]),
]

for i, (title, col, lines) in enumerate(card_data):
    x = 0.35 + i * 2.35
    add_card(slide, x, 1.5, 2.15, 5.3, CARD_LIGHT)
    add_textbox(slide, x + 0.1, 1.55, 1.95, 0.65, title, font_size=13, bold=True, color=col)
    add_hline(slide, x + 0.1, 2.2, 1.95, col)
    add_multiline(slide, x + 0.1, 2.35, 1.95, 4.2, lines, font_size=10, color=DARK_TEXT)


# =========================================================================
# SLIDE 4: Research Gap & Contribution
# =========================================================================
slide = add_slide(LY_BLANK)
slide_title(slide, "Research Gap & Our Contribution")

gap_data = [
    ["Aspect", "Prior Work", "Our Project"],
    ["Noise dimensions", "Study 1 dimension in isolation", "Combine prompt + paraphrase + item noise"],
    ["Item-level analysis", "Aggregate accuracy only", "Per-item noise score mapping"],
    ["Noise separation", "Not attempted", "Three-way variance decomposition"],
    ["Actionable filtering", "Report problems only", "Remove noisy items & measure impact"],
    ["Cross-source check", "Single paraphrase source", "Dual source (GPT-4o + Qwen)"],
    ["Model scale effect", "1-2 models", "4 models across 7B-72B"],
]
add_table(slide, 0.4, 1.5, 9.2, 2.6, len(gap_data), 3, gap_data, cell_font=11, header_font=12)

# Core RQ box
add_card(slide, 0.4, 4.4, 9.2, 0.8, BLUE_BG)
add_textbox(slide, 0.6, 4.45, 8.8, 0.65,
            "Core RQ: How much of the observed performance variation reflects\ntrue capability vs. evaluation noise?",
            font_size=15, bold=True, color=WHITE, alignment=PP_ALIGN.CENTER)

# Three sub-RQs
rqs = [
    ("RQ1", "How sensitive are rankings\nto prompt wording?", ACCENT1),
    ("RQ2", "Does paraphrasing questions\nchange accuracy & rank?", ACCENT2),
    ("RQ3", "Can item noise identification\nimprove reliability?", ACCENT3),
]
for i, (tag, text, col) in enumerate(rqs):
    x = 0.4 + i * 3.15
    add_card(slide, x, 5.5, 2.95, 1.5, CARD_LIGHT)
    add_textbox(slide, x + 0.1, 5.55, 0.7, 0.3, tag, font_size=14, bold=True, color=col)
    add_textbox(slide, x + 0.1, 5.85, 2.7, 0.9, text, font_size=12, color=DARK_TEXT)


# =========================================================================
# SLIDE 5: Experimental Setup
# =========================================================================
slide = add_slide(LY_BLANK)
slide_title(slide, "Experimental Setup: Models & Datasets")

# Models table
model_data = [
    ["Model", "Params", "Family", "Context"],
    ["LLaMA-3.1-8B", "8B", "LLaMA", "128K"],
    ["Qwen2.5-7B", "7B", "Qwen", "128K"],
    ["Qwen3-32B", "32B", "Qwen", "128K"],
    ["Qwen2.5-72B", "72B", "Qwen", "128K"],
]
add_textbox(slide, 0.5, 1.4, 4, 0.3, "Models (4 models, 7B–72B)", font_size=13, bold=True, color=ACCENT1)
add_table(slide, 0.4, 1.8, 4.5, 1.8, len(model_data), 4, model_data, cell_font=11, header_font=11)

# Dataset table
ds_data = [
    ["Benchmark", "Items", "Options", "Difficulty"],
    ["ARC-Challenge", "300", "4 (A-D)", "Medium"],
    ["MMLU-Pro", "300", "10 (A-J)", "Hard"],
]
add_textbox(slide, 5.3, 1.4, 4, 0.3, "Benchmarks", font_size=13, bold=True, color=ACCENT1)
add_table(slide, 5.2, 1.8, 4.4, 1.1, len(ds_data), 4, ds_data, cell_font=11, header_font=11)

add_card(slide, 5.2, 3.1, 4.4, 0.6, CARD_BLUE)
add_textbox(slide, 5.4, 3.15, 4.0, 0.5,
            "API: OpenRouter | Temp = 0.0 | Async batch",
            font_size=11, color=DARK_TEXT)

# Experiment overview
exp_data = [
    ["Experiment", "Focus", "Variables", "Evaluations"],
    ["Exp I", "Prompt wording noise", "18 variants × 4 models × 2 bench", "21,600"],
    ["Exp II", "Paraphrase sensitivity", "4 versions × 2 sources × 4 models", "4,800"],
    ["Exp III", "Item-level noise map", "Exp I + II → filter 10/20/30%", "Re-analysis"],
]
add_textbox(slide, 0.5, 4.1, 4, 0.3, "Experiment Overview", font_size=13, bold=True, color=ACCENT1)
add_table(slide, 0.4, 4.5, 9.2, 1.6, len(exp_data), 4, exp_data, cell_font=11, header_font=11)


# =========================================================================
# SLIDE 6: Exp I Design — Prompt Perturbation
# =========================================================================
slide = add_slide(LY_BLANK)
slide_title(slide, "Exp I Design: Prompt Perturbation (OFAT)")

# Prompt dimension table
pdim = [
    ["Dimension", "Level 0 (Base)", "Level 1", "Level 2"],
    ["Instruction", "Answer the following", "Choose the correct answer", "Select the best option"],
    ["Answer Fmt", "letter only", "letter + explanation", "full answer text"],
    ["Option Fmt", "A. text  B. text", "(A) text  (B) text", "1. text  2. text"],
    ["Framing", "(none)", "You are an expert", "—"],
]
add_table(slide, 0.4, 1.5, 9.2, 1.8, len(pdim), 4, pdim, cell_font=10, header_font=11)

# OFAT explanation
add_card(slide, 0.4, 3.5, 4.3, 3.6, CARD_LIGHT)
add_textbox(slide, 0.6, 3.55, 3.9, 0.3, "OFAT Design", font_size=13, bold=True, color=ACCENT1)
ofat = [
    "• 1 base variant (all at level 0)",
    "• 7 OFAT: change 1 dimension at a time",
    "   → isolate each dimension's effect",
    "• 10 random combos (from 54 possible)",
    "   → capture interaction effects",
    "• Total: 18 variants per model × dataset",
    "",
    "Metrics:",
    "  · Accuracy mean / std / range",
    "  · Item flip rate (%)",
    "  · Var(prompt) / Var(sampling) ratio",
    "  · Pairwise ranking reversal rate",
]
add_multiline(slide, 0.6, 3.95, 3.9, 3.0, ofat, font_size=10, color=DARK_TEXT)

# OFAT main effects figure
add_img(slide, os.path.join(FIG1, "fig2_ofat_main_effects.png"), 4.9, 3.5, width=4.8)


# =========================================================================
# SLIDE 7: Exp II & III Design
# =========================================================================
slide = add_slide(LY_BLANK)
slide_title(slide, "Exp II & III Design")

# Exp II
add_card(slide, 0.4, 1.5, 4.5, 2.8, CARD_LIGHT)
add_textbox(slide, 0.6, 1.55, 4.1, 0.3, "Exp II: Paraphrase Sensitivity", font_size=14, bold=True, color=ACCENT2)
exp2_lines = [
    "• Each question → 3 meaning-preserving paraphrases",
    "• Dual source: GPT-4o + Qwen2.5-72B",
    "• Subset: 150 shared items per benchmark",
    "• 4 versions/item: original + 3 paraphrases",
    "• Parse failure → count as incorrect (primary)",
    "• Clean subset analysis (supplemental check)",
]
add_multiline(slide, 0.6, 2.0, 4.1, 2.0, exp2_lines, font_size=11, color=DARK_TEXT)

# Exp III
add_card(slide, 0.4, 4.55, 4.5, 2.6, CARD_LIGHT)
add_textbox(slide, 0.6, 4.6, 4.1, 0.3, "Exp III: Item Noise Score & Filtering", font_size=14, bold=True, color=ACCENT3)
exp3_lines = [
    "Noise(q) = 1 − |2c(q) − N(q)| / N(q)",
    "  c(q) = #correct across prompts/paraphrases",
    "  All-correct or all-wrong → noise = 0",
    "  50/50 → noise ≈ 1 (maximally unstable)",
    "",
    "Filter top 10% / 20% / 30% noisiest items",
    "→ Re-evaluate all metrics on clean subset",
]
add_multiline(slide, 0.6, 5.1, 4.1, 2.0, exp3_lines, font_size=11, color=DARK_TEXT)

# Flow diagram
add_card(slide, 5.2, 1.5, 4.4, 5.65, CARD_BLUE)
add_textbox(slide, 5.4, 1.55, 4.0, 0.3, "Three-Experiment Flow", font_size=14, bold=True, color=BLUE_BG)
flow = [
    "┌──────────────────────────┐",
    "│  Exp I: Prompt Perturb.  │",
    "│  300 items × 18 variants │",
    "│  → Acc, flip rate, VarR  │",
    "└────────────┬─────────────┘",
    "             │ item noise",
    "             ▼",
    "┌──────────────────────────┐",
    "│  Exp II: Paraphrase      │",
    "│  150 items × 4 versions  │",
    "│  → Cross-source check    │",
    "└────────────┬─────────────┘",
    "             │ combined",
    "             ▼",
    "┌──────────────────────────┐",
    "│  Exp III: Noise Map      │",
    "│  Noise(q) → filter noisy │",
    "│  → Re-evaluate stability │",
    "└──────────────────────────┘",
]
add_multiline(slide, 5.4, 2.0, 4.0, 5.0, flow, font_size=10, color=DARK_TEXT, font_name="Courier New")


# =========================================================================
# SLIDE 8: Exp I Results — Accuracy
# =========================================================================
slide = add_slide(LY_BLANK)
slide_title(slide, "Exp I: Prompt Sensitivity — Accuracy")

add_img(slide, os.path.join(FIG1, "fig1_accuracy_per_variant.png"), 0.2, 1.4, width=6.8)

add_card(slide, 7.2, 1.4, 2.5, 5.5, CARD_LIGHT)
add_textbox(slide, 7.35, 1.45, 2.2, 0.3, "Key Findings", font_size=14, bold=True, color=ACCENT1)

f_lines = [
    "Accuracy Range (ARC):",
    "  72B: 33.3%",
    "  32B: 14.7%",
    "  7B:   20.3%",
    "  8B:   10.3%",
    "",
    "Range (MMLU):",
    "  72B: 21.4%",
    "  32B: 13.1%",
    "  8B:   12.4%",
    "",
    "MMLU Item Flip Rate:",
    "  All models > 62%",
    "",
    "⚠ Bigger ≠ Stable",
    "72B has best acc but",
    "worst range (33%)",
]
add_multiline(slide, 7.35, 1.85, 2.2, 4.8, f_lines, font_size=10, color=DARK_TEXT)


# =========================================================================
# SLIDE 9: Exp I — Variance & Scale
# =========================================================================
slide = add_slide(LY_BLANK)
slide_title(slide, "Exp I: Variance Decomposition & Scale")

add_img(slide, os.path.join(FIG1, "fig3_variance_decomposition.png"), 0.1, 1.4, width=4.8)
add_img(slide, os.path.join(FIG1, "fig5_scale_analysis.png"), 5.0, 1.4, width=4.8)

# Summary table
t1 = [
    ["Model", "ARC Acc", "ARC Range", "ARC VarR", "MMLU Acc", "MMLU Range", "MMLU VarR"],
    ["LLaMA-8B", "79.2%", "10.3%", "2.3×", "39.9%", "12.4%", "1.7×"],
    ["Qwen-7B", "86.9%", "20.3%", "7.4×", "45.3%", "10.3%", "1.1×"],
    ["Qwen-32B", "91.7%", "14.7%", "11.1×", "58.2%", "13.1%", "1.4×"],
    ["Qwen-72B", "89.6%", "33.3%", "49.1×", "54.5%", "21.4%", "3.4×"],
]
add_table(slide, 0.2, 5.0, 9.6, 1.7, len(t1), 7, t1, cell_font=10, header_font=10)
add_textbox(slide, 0.2, 6.7, 9.6, 0.3,
            "VarR = Var(prompt)/Var(sampling). Values > 1 → prompt noise dominates. 72B: VarR = 49× on ARC!",
            font_size=10, color=MID_GRAY, alignment=PP_ALIGN.CENTER)


# =========================================================================
# SLIDE 10: Exp I — Ranking Reversals
# =========================================================================
slide = add_slide(LY_BLANK)
slide_title(slide, "Exp I: Ranking Reversals & Categories")

add_img(slide, os.path.join(FIG1, "fig8_reversal_summary.png"), 0.1, 1.4, width=4.8)
add_img(slide, os.path.join(FIG1, "fig7_category_heatmap.png"), 5.0, 1.4, width=4.8)

# Interpretations
add_card(slide, 0.3, 4.8, 4.5, 2.3, CARD_LIGHT)
add_textbox(slide, 0.5, 4.85, 4.1, 0.3, "Ranking Is Prompt-Dependent", font_size=13, bold=True, color=RED_ALERT)
r_lines = [
    "• 32B vs 72B reverses 67% on ARC",
    "• Just changing the prompt flips \"better\"",
    "• MMLU more stable (wider absolute gaps)",
    "• \"Model A > B\" on one prompt ≠ on another",
]
add_multiline(slide, 0.5, 5.25, 4.1, 1.5, r_lines, font_size=11, color=DARK_TEXT)

add_card(slide, 5.1, 4.8, 4.5, 2.3, CARD_BLUE)
add_textbox(slide, 5.3, 4.85, 4.1, 0.3, "Category Sensitivity", font_size=13, bold=True, color=BLUE_BG)
c_lines = [
    "• Uneven sensitivity across MMLU subjects",
    "• Biology, Psychology most prompt-sensitive",
    "• Math, CS relatively more stable",
    "• → Noise is content-dependent",
]
add_multiline(slide, 5.3, 5.25, 4.1, 1.5, c_lines, font_size=11, color=DARK_TEXT)


# =========================================================================
# SLIDE 11: Exp II — Accuracy Stability
# =========================================================================
slide = add_slide(LY_BLANK)
slide_title(slide, "Exp II: Paraphrase — Accuracy Stability")

add_img(slide, os.path.join(FIG2, "fig1_accuracy_by_version.png"), 0.1, 1.4, width=6.5)

add_card(slide, 6.8, 1.4, 2.9, 5.5, CARD_LIGHT)
add_textbox(slide, 7.0, 1.45, 2.5, 0.3, "Key Findings", font_size=14, bold=True, color=ACCENT2)
e2_f = [
    "Paraphrase noise is real",
    "but weaker than prompt",
    "noise.",
    "",
    "ARC accuracy stability:",
    "  Large (32B,72B): ±1-2%",
    "  Small (7B, 8B):  ±3-5%",
    "",
    "MMLU (harder task):",
    "  More variation observed",
    "  Qwen2.5-7B: up to 5%",
    "  swing per paraphrase",
    "",
    "Both GPT-4o and Qwen",
    "sources yield similar",
    "patterns → source bias",
    "is minimal",
]
add_multiline(slide, 7.0, 1.85, 2.5, 4.8, e2_f, font_size=10, color=DARK_TEXT)


# =========================================================================
# SLIDE 12: Exp II — Flip Rate & Ranking
# =========================================================================
slide = add_slide(LY_BLANK)
slide_title(slide, "Exp II: Flip Rate & Ranking Stability")

add_img(slide, os.path.join(FIG2, "fig2_flip_rate.png"), 0.1, 1.4, width=4.8)
add_img(slide, os.path.join(FIG2, "fig6_cross_source_accuracy.png"), 5.0, 1.4, width=4.8)

add_img(slide, os.path.join(FIG2, "fig4_rank_distribution.png"), 0.1, 4.1, width=4.8, height=2.5)
add_img(slide, os.path.join(FIG2, "fig7_cross_source_significance.png"), 5.0, 4.1, width=4.8, height=2.5)


# =========================================================================
# SLIDE 13: Exp II — Variance & Cross-Experiment
# =========================================================================
slide = add_slide(LY_BLANK)
slide_title(slide, "Exp II: Variance & Cross-Experiment Comparison")

add_img(slide, os.path.join(FIG2, "fig9_variance_decomposition.png"), 0.1, 1.3, width=6.0)
add_img(slide, os.path.join(FIG2, "fig8_cross_experiment_flip.png"), 0.1, 4.3, width=6.0)

add_card(slide, 6.3, 1.3, 3.4, 5.7, CARD_LIGHT)
add_textbox(slide, 6.5, 1.35, 3.0, 0.3, "Key Insights", font_size=14, bold=True, color=ACCENT2)
vi = [
    "Variance Decomposition:",
    "  Prompt var (red) dominates",
    "  for most model-dataset pairs",
    "  Test-set var (blue) is small",
    "  but non-negligible (2-48%)",
    "",
    "Cross-Experiment:",
    "  Exp I flip rates are 2-12×",
    "  higher than Exp II",
    "",
    "  ARC best case:",
    "  Exp I: 35% → Exp II: 4%",
    "  (Qwen2.5-72B)",
    "",
    "  MMLU worst case:",
    "  Exp I: 67% → Exp II: 29%",
    "  (Qwen2.5-7B)",
    "",
    "→ Prompt perturbation is",
    "  the dominant noise source",
]
add_multiline(slide, 6.5, 1.75, 3.0, 5.0, vi, font_size=10, color=DARK_TEXT)


# =========================================================================
# SLIDE 14: Exp III — Noise Distribution
# =========================================================================
slide = add_slide(LY_BLANK)
slide_title(slide, "Exp III: Item-Level Noise Distribution")

add_img(slide, os.path.join(FIG3, "fig1_noise_distribution.png"), 0.1, 1.4, width=6.5)

add_card(slide, 6.8, 1.4, 2.9, 2.5, CARD_LIGHT)
add_textbox(slide, 7.0, 1.45, 2.5, 0.3, "Noise Stats", font_size=13, bold=True, color=ACCENT3)
ns = [
    "ARC-Challenge:",
    "  Mean: 0.235 | Med: 0.077",
    "  Items > 0.5: 22.0%",
    "",
    "MMLU-Pro:",
    "  Mean: 0.452 | Med: 0.443",
    "  Items > 0.5: 39.3%",
    "",
    "MMLU ≈ 2× noisier than ARC",
]
add_multiline(slide, 7.0, 1.85, 2.5, 2.0, ns, font_size=10, color=DARK_TEXT)

add_img(slide, os.path.join(FIG3, "fig8_noise_vs_difficulty.png"), 0.1, 4.3, width=4.2)
add_img(slide, os.path.join(FIG3, "fig9_category_noise.png"), 4.5, 4.3, width=3.0)

add_card(slide, 7.6, 4.3, 2.1, 2.7, CARD_BLUE)
add_textbox(slide, 7.7, 4.35, 1.9, 0.2, "Noise vs Difficulty", font_size=11, bold=True, color=BLUE_BG)
nd = [
    "ARC: r = 0.81",
    "  Harder → noisier",
    "  (predictable)",
    "MMLU: r ≈ −0.05",
    "  No correlation!",
    "  (harder to fix)",
]
add_multiline(slide, 7.7, 4.65, 1.9, 2.0, nd, font_size=10, color=DARK_TEXT)


# =========================================================================
# SLIDE 15: Exp III — Noise Correlation
# =========================================================================
slide = add_slide(LY_BLANK)
slide_title(slide, "Exp III: Noise Correlation & Properties")

add_img(slide, os.path.join(FIG3, "fig7_noise_correlation.png"), 0.1, 1.4, width=4.5)
add_img(slide, os.path.join(FIG3, "fig4_variance_ratio.png"), 4.8, 1.4, width=5.0)

add_card(slide, 0.3, 4.7, 4.3, 2.4, CARD_LIGHT)
add_textbox(slide, 0.5, 4.75, 3.9, 0.3, "Correlation Patterns", font_size=13, bold=True, color=ACCENT3)
cp = [
    "• Moderate cross-model r ≈ 0.3-0.5",
    "  → Some items universally noisy",
    "• Higher within-family correlation",
    "• Exp I noise ⊥ Exp II noise",
    "  → Independent noise dimensions",
]
add_multiline(slide, 0.5, 5.15, 3.9, 1.5, cp, font_size=11, color=DARK_TEXT)

add_card(slide, 4.9, 4.7, 4.8, 2.4, CARD_BLUE)
add_textbox(slide, 5.1, 4.75, 4.4, 0.3, "Variance Ratio Under Filtering", font_size=13, bold=True, color=BLUE_BG)
vr = [
    "• Removing noisy items reduces VarR",
    "  on ARC but not always on MMLU",
    "• ARC: noise concentrated in few items",
    "  → filtering is effective",
    "• MMLU: noise is diffuse → partial effect",
]
add_multiline(slide, 5.1, 5.15, 4.4, 1.5, vr, font_size=11, color=DARK_TEXT)


# =========================================================================
# SLIDE 16: Exp III — Filtering Results
# =========================================================================
slide = add_slide(LY_BLANK)
slide_title(slide, "Exp III: Noise Removal — Filtering Effect")

add_img(slide, os.path.join(FIG3, "fig3_flip_rate_reduction.png"), 0.1, 1.4, width=6.0)
add_img(slide, os.path.join(FIG3, "fig2_exp1_stability.png"), 0.1, 4.0, width=4.5)

# Filtering table
ft = [
    ["Dataset", "Model", "Baseline", "−10%", "−20%", "−30%", "Δ(30%)"],
    ["ARC", "LLaMA-8B", "34.0%", "29.6%", "26.7%", "19.0%", "−44%"],
    ["ARC", "Qwen-7B", "36.7%", "31.9%", "27.5%", "24.8%", "−32%"],
    ["ARC", "Qwen-32B", "30.0%", "25.9%", "21.7%", "19.0%", "−37%"],
    ["ARC", "Qwen-72B", "46.7%", "43.7%", "42.5%", "40.0%", "−14%"],
    ["MMLU", "LLaMA-8B", "64.0%", "62.2%", "60.8%", "57.1%", "−11%"],
    ["MMLU", "Qwen-7B", "61.3%", "57.8%", "55.0%", "49.5%", "−19%"],
    ["MMLU", "Qwen-32B", "65.3%", "65.9%", "63.3%", "62.9%", "−4%"],
    ["MMLU", "Qwen-72B", "60.7%", "57.8%", "53.3%", "50.5%", "−17%"],
]
add_table(slide, 4.7, 4.0, 5.1, 3.0, len(ft), 7, ft, cell_font=9, header_font=10)


# =========================================================================
# SLIDE 17: Exp III — Three-Way Variance & Paraphrase Stability
# =========================================================================
slide = add_slide(LY_BLANK)
slide_title(slide, "Exp III: Three-Way Variance & Reversal Reduction")

add_img(slide, os.path.join(FIG3, "fig6_three_way_variance.png"), 0.1, 1.4, width=6.0)
add_img(slide, os.path.join(FIG3, "fig5_reversal_reduction.png"), 0.1, 4.3, width=4.8)
add_img(slide, os.path.join(FIG3, "fig10_exp2_stability.png"), 5.1, 4.3, width=4.8)

add_card(slide, 6.3, 1.4, 3.4, 2.5, CARD_LIGHT)
add_textbox(slide, 6.5, 1.45, 3.0, 0.3, "Summary", font_size=13, bold=True, color=ACCENT3)
sm = [
    "Three-Way Variance:",
    "  ARC: prompt = 51-98%",
    "  MMLU: 33-69% prompt",
    "",
    "After −30% filter:",
    "  ✓ ARC flip ↓ 14-44%",
    "  ✓ ARC reversals ↓",
    "  △ MMLU flip ↓ 4-19%",
    "  △ MMLU reversal: limited",
]
add_multiline(slide, 6.5, 1.85, 3.0, 2.0, sm, font_size=10, color=DARK_TEXT)


# =========================================================================
# SLIDE 18: Summary Dashboard
# =========================================================================
slide = add_slide(LY_BLANK)
slide_title(slide, "Summary Dashboards")

add_img(slide, os.path.join(FIG1, "fig9_summary_dashboard.png"), 0.05, 1.3, width=4.9)
add_img(slide, os.path.join(FIG3, "fig11_summary_dashboard.png"), 5.05, 1.3, width=4.9)

add_textbox(slide, 0.05, 6.5, 4.9, 0.4, "Exp I: Prompt Perturbation",
            font_size=14, bold=True, color=ACCENT1, alignment=PP_ALIGN.CENTER)
add_textbox(slide, 5.05, 6.5, 4.9, 0.4, "Exp III: High-Noise Item Analysis",
            font_size=14, bold=True, color=ACCENT3, alignment=PP_ALIGN.CENTER)


# =========================================================================
# SLIDE 19: Discussion — Key Takeaways
# =========================================================================
slide = add_slide(LY_BLANK)
slide_title(slide, "Discussion & Key Takeaways")

find_data = [
    ("① Noise is Pervasive", RED_ALERT,
     ["Prompt → 10-33% acc swing",
      "62-67% MMLU items flip",
      "Single scores unreliable"]),
    ("② Prompt >> Paraphrase", ACCENT2,
     ["Prompt noise 2-12× stronger",
      "Standardizing prompts is",
      "the highest-leverage fix"]),
    ("③ Bigger ≠ More Stable", ACCENT1,
     ["72B: best acc, worst range",
      "(33%, VarR=49×). Accuracy",
      "and robustness ≠ same"]),
    ("④ MMLU >> ARC in Noise", RED_ALERT,
     ["MMLU noise: 0.452 vs 0.235",
      "More options, harder content",
      "→ more evaluation noise"]),
    ("⑤ Filtering Helps Partly", ACCENT3,
     ["ARC flip ↓ up to 44%",
      "MMLU flip ↓ only 4-19%",
      "Benchmark-dependent"]),
    ("⑥ Rankings Are Fragile", ACCENT2,
     ["32B vs 72B reverses 67%",
      "on ARC. One prompt  can't",
      "determine who is better"]),
]

for i, (title, col, lines) in enumerate(find_data):
    r, c = divmod(i, 3)
    x = 0.35 + c * 3.15
    y = 1.4 + r * 3.05
    add_card(slide, x, y, 2.95, 2.75, CARD_LIGHT)
    add_textbox(slide, x + 0.1, y + 0.05, 2.75, 0.35, title, font_size=13, bold=True, color=col)
    add_hline(slide, x + 0.1, y + 0.4, 2.75, col)
    add_multiline(slide, x + 0.1, y + 0.55, 2.75, 2.0, lines, font_size=11, color=DARK_TEXT)


# =========================================================================
# SLIDE 20: Limitations & Future Work
# =========================================================================
slide = add_slide(LY_BLANK)
slide_title(slide, "Limitations & Future Directions")

# Limitations
add_card(slide, 0.4, 1.4, 4.4, 5.6, CARD_LIGHT)
add_textbox(slide, 0.6, 1.45, 4.0, 0.3, "Current Limitations", font_size=15, bold=True, color=RED_ALERT)
lim = [
    "① Paraphrase fidelity unverified",
    "   No check that meaning is exactly",
    "   preserved → possible semantic drift",
    "",
    "② No repeated runs at temp=0",
    "   Cannot separate API randomness",
    "   from genuine noise effects",
    "",
    "③ Scale constraints",
    "   4 models (3 Qwen), 2 benchmarks,",
    "   18 variants (vs literature 50-4800)",
    "",
    "④ Non-orthogonal OFAT design",
    "   Main effects may confound with",
    "   interaction terms",
    "",
    "⑤ Data contamination unassessed",
    "   ARC/MMLU likely in training data",
]
add_multiline(slide, 0.6, 1.85, 4.0, 4.8, lim, font_size=11, color=DARK_TEXT)

# Future
add_card(slide, 5.1, 1.4, 4.5, 5.6, CARD_BLUE)
add_textbox(slide, 5.3, 1.45, 4.1, 0.3, "Future Directions", font_size=15, bold=True, color=ACCENT3)
fut = [
    "① Semantic verification",
    "   GPT-4o scoring + human spot-check",
    "   for paraphrase quality",
    "",
    "② Repeated trials (3-5×)",
    "   Establish baseline variance for",
    "   temp=0 API calls",
    "",
    "③ Diverse model families",
    "   Add Mistral, Gemma, GPT to test",
    "   cross-family generalizability",
    "",
    "④ IRT-based item analysis",
    "   Separate difficulty from instability",
    "   (noise ≠ hardness)",
    "",
    "⑤ Robustness certification",
    "   Min items to remove for reversal",
    "   → benchmark cleaning tool",
]
add_multiline(slide, 5.3, 1.85, 4.1, 4.8, fut, font_size=11, color=DARK_TEXT)


# =========================================================================
# SLIDE 21: Conclusion & Thank You
# =========================================================================
slide = add_slide(LY_BLANK)
set_bg(slide, BLUE_BG)

add_textbox(slide, 0.6, 1.0, 8.8, 0.6, "Conclusion",
            font_size=34, bold=True, color=WHITE, font_name="Calibri Light",
            alignment=PP_ALIGN.CENTER)
add_hline(slide, 3.5, 1.65, 3.0, ACCENT2)

conc = [
    "✦  Benchmark scores carry substantial evaluation noise — single scores unreliable",
    "",
    "✦  Prompt perturbation is the dominant source (item flip > 62%, VarR up to 49×)",
    "",
    "✦  Larger models are NOT more robust (72B: acc=90% but range=33%)",
    "",
    "✦  Item-level noise mapping → targeted cleaning improves stability up to 44%",
    "",
    "✦  Paraphrase noise is 2-12× weaker; cross-source consistency = 100%",
]
add_multiline(slide, 0.6, 1.9, 8.8, 3.5, conc, font_size=14, color=LIGHT_GRAY,
              alignment=PP_ALIGN.CENTER)

add_textbox(slide, 0.6, 5.2, 8.8, 0.7, "Thank You!   Q & A",
            font_size=32, bold=True, color=WHITE, font_name="Calibri Light",
            alignment=PP_ALIGN.CENTER)
add_textbox(slide, 0.6, 6.0, 8.8, 0.4, "ST5230 Applied NLP — Group 1",
            font_size=14, color=RGBColor(0xA0, 0xC0, 0xE0), alignment=PP_ALIGN.CENTER)


# =========================================================================
# Save
# =========================================================================
out = os.path.join(BASE, "presentation_template.pptx")
prs.save(out)
print(f"✅ Saved to {out}")
print(f"   Total slides: {len(prs.slides)}")
